from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(output_path):
    prs = Presentation()
    
    # helper function to set slide background to dark
    def set_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(3, 7, 18)

    # --- SLIDE 1: TITLE ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_dark_bg(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SMARTGUARD AI"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 245, 255)
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Enterprise-Grade Threat Intelligence Pipeline\nComprehensive Technical Presentation"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # --- SLIDE 2: TECH STACK RATIONALE ---
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_bg(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = "1. Technical Stack Rationale"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 245, 255)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = "BACKEND: FastAPI (Python)"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "- Chosen for high-performance async processing of heavy file scans."
    p.font.color.rgb = RGBColor(148, 163, 184)
    
    p = tf.add_paragraph()
    p.text = "FRONTEND: Streamlit (Obsidian Theme)"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "- Rapid delivery of reactive data components and SOC-style visuals."
    p.font.color.rgb = RGBColor(148, 163, 184)

    # --- SLIDE 3: DATA LIFECYCLE ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_bg(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = "2. Data Lifecycle & Privacy"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 245, 255)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = "INGESTION: Secure Byte-Stream via TLS 1.3."
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "PROCESSING: Zero-Retention in-memory decomposition."
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "STORAGE: Metadata only (SHA-256, results) in isolated JSON history."
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "PRIVACY: Session isolation ensures users only see their own archives."
    p.font.color.rgb = RGBColor(0, 245, 255)

    # --- SLIDE 4: MACHINE LEARNING (NSL-KDD) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_bg(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = "3. ML Core: The NSL-KDD Advantage"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 245, 255)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = "DATASET: NSL-KDD from UNB Canadian Institute for Cybersecurity."
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "TRAINING: 125,000+ records across four attack categories."
    p.font.color.rgb = RGBColor(148, 163, 184)
    
    p = tf.add_paragraph()
    p.text = "ARCHITECTURE: RF & Neural Ensemble for pattern recognition."
    p.font.color.rgb = RGBColor(148, 163, 184)

    # --- SLIDE 5: ELITE FEATURES DEEP-DIVE ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_bg(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = "4. Elite Mechanical Features"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 245, 255)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = "CINEMATIC SCANS: Visual neural bridge feedback during deep analysis."
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "REMEDIATION: Risk-aware checklists (Isolation, Quarantine, Recovery)."
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "FORENSIC COMPARE: High-contrast pattern analysis between specimens."
    p.font.color.rgb = RGBColor(255, 255, 255)

    # --- SLIDE 6: SUMMARY ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    set_dark_bg(slide)
    
    title_shape = slide.shapes.title
    title_shape.text = "Conclusion: SmartGuard AI"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 245, 255)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Securing the Digital World with Hybrid Intelligence."
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    p = tf.add_paragraph()
    p.text = "- Clinical Detection Accuracy"
    p.font.color.rgb = RGBColor(148, 163, 184)
    p = tf.add_paragraph()
    p.text = "- Cinematic User Experience"
    p.font.color.rgb = RGBColor(148, 163, 184)
    p = tf.add_paragraph()
    p.text = "- Enterprise Privacy Standards"
    p.font.color.rgb = RGBColor(148, 163, 184)

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "SmartGuard_AI_Presentation.pptx"
    create_presentation(out)
